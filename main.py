import json
from slides import Slide
import os
from pptx import Presentation
from pptx.util import Inches
import numpy as np
import matplotlib.pyplot as plt

def create_title_slide(slide_layout, title, content):
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    if content:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

def create_text_slide(slide_layout, title, content):
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    if content:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

def create_list_slide(slide_layout, title, content):
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    if content:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = ''
        for item in content:
            content_placeholder.text += f"- {item['text']}\n"

def create_picture_slide(slide_layout, title, image_filename):
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    left = Inches(2)
    top = Inches(2)
    pic = slide.shapes.add_picture(image_filename, left, top, width=Inches(6))

def create_plot_slide(slide_layout, title, data_filename, config):
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    if len(slide.placeholders) > 1:
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = ''

    # Load data from the .dat file using numpy
    data = np.loadtxt(data_filename, delimiter=';')

    # Assuming the data contains two columns, x and y
    x = data[:, 0]
    y = data[:, 1]

    # Create a plot using matplotlib
    plt.plot(x, y)
    plt.xlabel(config.get('x-label', ''))
    plt.ylabel(config.get('y-label', ''))

    # Save the plot to a file
    plot_filename = 'plot.png'
    plt.savefig(plot_filename)

    # Add the plot to the slide
    slide.shapes.add_picture(plot_filename, Inches(2), Inches(2), width=Inches(6))


# Asking for the configuration file from the user
if __name__ == "__main__":
    json_file_path = input("Enter the path to the configuration file: ")

    # Load JSON data from the specified file
    with open(json_file_path) as json_file:
        json_data = json.load(json_file)

    # Create a presentation object
    presentation = Presentation()

    # Define slide layouts
    slide_layouts = presentation.slide_layouts

    # Process each slide data from the JSON
    for slide_data in json_data['presentation']:
        slide_type = slide_data.get('type')
        slide_title = slide_data.get('title')
        slide_content = slide_data.get('content')

        if slide_type == 'title':
            create_title_slide(slide_layouts[0], slide_title, slide_content)

        elif slide_type == 'text':
            create_text_slide(slide_layouts[1], slide_title, slide_content)

        elif slide_type == 'list':
            create_list_slide(slide_layouts[1], slide_title, slide_content)

        elif slide_type == 'picture':
            image_filename = os.path.join(os.path.dirname(json_file_path), slide_content)
            create_picture_slide(slide_layouts[5], slide_title, image_filename)

        elif slide_type == 'plot':
            data_filename = os.path.join(os.path.dirname(json_file_path), slide_content)
            create_plot_slide(slide_layouts[5], slide_title, data_filename, slide_data.get('configuration', {}))

    # Save the presentation to a file
    presentation.save('output_presentation.pptx')